
'''
- 실행 방법: python extract_ppt.py "PPT 파일 경로"
- 예시: python extract_ppt.py "투자 성향 진단 테스트(성인용).pptx"
'''
import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text_and_images_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    
    # 결과 저장 폴더 생성
    output_folder = os.path.splitext(os.path.basename(pptx_path))[0] + "_extracted"
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    print(f"--- 텍스트 추출 ({pptx_path}) ---\n")
    slide_count = 0
    for slide in prs.slides:
        slide_count += 1
        print(f"--- Slide {slide_count} ---
")
        
        text_runs = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
        if text_runs:
            print("\n".join(text_runs))
            print("-" * 20)

    print(f"\n--- 이미지 추출 ({pptx_path}) ---\n")
    image_count = 0
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                image_bytes = image.blob
                
                # 이미지 파일명 생성
                image_ext = image.ext
                image_filename = f"slide_{i+1}_image_{image_count}.{image_ext}"
                image_path = os.path.join(output_folder, image_filename)
                
                # 이미지 저장
                with open(image_path, "wb") as f:
                    f.write(image_bytes)
                
                print(f"이미지 저장: {image_path}")
                image_count += 1

    print(f"\n추출 완료. ''{output_folder}'' 폴더를 확인하세요.")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("사용법: python extract_ppt.py \"path/to/your/presentation.pptx\"")
        sys.exit(1)
    
    file_path = sys.argv[1]
    if not os.path.exists(file_path):
        print(f"오류: 파일 ''{file_path}''를 찾을 수 없습니다.")
        sys.exit(1)

    extract_text_and_images_from_pptx(file_path)
